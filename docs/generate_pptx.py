from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ACCENT = RGBColor(0x3B, 0x5B, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x5A, 0x5F, 0x6D)
GREEN = RGBColor(0x2B, 0x8A, 0x3E)
RED = RGBColor(0xC9, 0x2A, 0x2A)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xF8)
ACCENT_LIGHT = RGBColor(0xED, 0xF2, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=MUTED, bold_prefix=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        if bold_prefix and ": " in item:
            parts = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = parts[0] + ": "
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = TEXT
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
        p.level = 0
    return txBox

def add_table(slide, left, top, width, height, headers, rows, font_size=11):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BG

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = MUTED
                paragraph.font.name = "Calibri"

    return table_shape

def add_accent_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_box(slide, left, top, width, height, title, subtitle, fill_color, text_color=WHITE, border_color=None, font_size_title=14, font_size_sub=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(font_size_title)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(font_size_sub)
        run2.font.color.rgb = text_color
        run2.font.name = "Calibri"
    return shape

def add_arrow(slide, start_left, start_top, end_left, end_top, color=MUTED):
    connector = slide.shapes.add_connector(
        1, Inches(start_left), Inches(start_top), Inches(end_left), Inches(end_top)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

def add_metric_box(slide, left, top, value, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2.5), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE2, 0xE5, 0xE9)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = value
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(11)
    run2.font.color.rgb = MUTED
    run2.font.name = "Calibri"


# ===== SLIDE 1: TITLE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF8, 0xF9, 0xFF))
add_accent_bar(slide)
add_text_box(slide, 1.5, 0.8, 10, 0.5, "ENGR 5785G - Real-Time Data Analytics for IoT", font_size=14, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.8, 11, 1.5, "Real-Time Cheat Detection\nfor Multiplayer Games", font_size=40, bold=True, color=TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.8, 10, 1.2, "A server-side streaming analytics pipeline using Kafka, Go, XGBoost,\nand LSTM Autoencoder to detect cheats in under 5 seconds\nwith full Grafana observability.", font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.5, 10, 0.5, "Go 1.22  |  Python 3.11  |  Apache Kafka  |  XGBoost + LSTM  |  Real-Time 60Hz  |  Docker Compose", font_size=13, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.5, 10, 0.5, "Ontario Tech University | Summer 2026", font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)


# ===== SLIDE 2: PROBLEM + IoT CONTEXT =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Problem Definition & IoT Context", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "Why real-time cheat detection matters and how it maps to IoT analytics", font_size=16, color=MUTED)

add_text_box(slide, 0.8, 1.6, 3.8, 0.4, "The Challenge", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 0.8, 2.0, 3.8, 2.5, [
    "Cheating costs the gaming industry billions annually",
    "Client-side anti-cheat (VAC, EAC) is bypassable",
    "Post-match analysis takes hours to days",
    "Client kernel drivers raise privacy concerns",
], font_size=13)

add_text_box(slide, 5, 1.6, 3.8, 0.4, "Our Approach", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 5, 2.0, 3.8, 2.5, [
    "Server-side: behavioral detection, no client trust",
    "Real-time: streaming pipeline via Apache Kafka",
    "Dual-paradigm ML: supervised + unsupervised",
    "Rule engine: deterministic, zero false positives",
    "Detection in under 5 seconds",
], font_size=13)

add_text_box(slide, 9.2, 1.6, 3.8, 0.4, "IoT Parallels", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 9.2, 2.0, 3.8, 2.5, [
    "Sensors: 60Hz telemetry like industrial IoT",
    "Edge Processing: sliding-window aggregation",
    "Stream Processing: Kafka (~540 events/sec)",
    "Anomaly Detection: LSTM autoencoder (same as IoT sensor anomaly detection)",
], font_size=13)

add_text_box(slide, 0.8, 4.7, 12, 0.4, "Related Work", font_size=18, bold=True, color=TEXT)
add_table(slide, 0.8, 5.2, 11.5, 1.8,
    ["Approach", "Examples", "Strengths", "Limitations"],
    [
        ["Client-Side Kernel", "VAC, EAC, BattlEye, Vanguard", "Detects known cheat software", "Kernel access; bypassable; privacy"],
        ["Statistical Post-Match", "CSGO Overwatch, manual review", "Human judgment, low FP", "Hours/days delay; doesn't scale"],
        ["ML-Based Behavioral", "Research papers (CS2CD)", "Detects behavioral patterns", "Usually offline; single-paradigm"],
    ], font_size=11)


# ===== SLIDE 3: ARCHITECTURE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "System Architecture", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "7 Docker containers, 4 Kafka topics, 3 detection layers", font_size=16, color=MUTED)

BLUE = RGBColor(0x3B, 0x5B, 0xDB)
TEAL = RGBColor(0x0C, 0x85, 0x99)
ORANGE = RGBColor(0xE6, 0x77, 0x00)
PURPLE = RGBColor(0x70, 0x48, 0xE8)
DARK_GREEN = RGBColor(0x2B, 0x8A, 0x3E)

add_box(slide, 0.5, 1.7, 2.4, 1.2, "Browser (:80)", "HTML5 Canvas\nWASD + Mouse\nF5-F8 Cheats", BLUE)
add_box(slide, 4.0, 1.7, 3.5, 1.2, "Game Server (:8080)", "Authoritative Game Loop\n60Hz Tick | 8 Bots + Humans\nPhysics | Raycast | Spatial", TEAL)
add_box(slide, 9.0, 1.7, 3.5, 1.2, "Apache Kafka (:9092)", "telemetry.raw | events.kills\nfeatures.computed\nalerts.detections", PURPLE)

add_box(slide, 4.0, 3.5, 3.5, 1.0, "Feature Engine (Go)", "Embedded in Game Server\n3 Windows (1s, 5s, 30s) | 18 Features | 6 Rules", TEAL)

add_box(slide, 4.0, 5.1, 3.5, 1.0, "ML Service (:8000)", "Python | XGBoost (per-tick)\nLSTM Autoencoder (per-sec)", ORANGE, text_color=WHITE)

add_box(slide, 9.0, 3.5, 1.6, 1.0, "Prometheus", "(:9091)\nScrapes every 5s", RGBColor(0x86, 0x4A, 0x15), text_color=WHITE, font_size_title=12, font_size_sub=9)
add_box(slide, 11.0, 3.5, 1.5, 1.0, "Grafana", "(:3000)\n2 Dashboards", DARK_GREEN, text_color=WHITE, font_size_title=12, font_size_sub=9)

add_arrow(slide, 2.9, 2.3, 4.0, 2.3, MUTED)
add_text_box(slide, 3.0, 1.95, 0.8, 0.3, "WS", font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

add_arrow(slide, 7.5, 2.3, 9.0, 2.3, MUTED)
add_arrow(slide, 5.75, 2.9, 5.75, 3.5, MUTED)
add_arrow(slide, 7.5, 3.9, 9.0, 3.9, MUTED)
add_arrow(slide, 5.75, 4.5, 5.75, 5.1, MUTED)
add_arrow(slide, 9.0, 5.4, 7.5, 5.4, MUTED)
add_arrow(slide, 7.5, 5.6, 9.0, 5.6, MUTED)
add_arrow(slide, 10.6, 3.5, 10.6, 2.9, MUTED)
add_arrow(slide, 11.0, 4.0, 10.6, 4.0, MUTED)

add_box(slide, 9.0, 5.1, 3.5, 0.7, "Kafka Topics (4 partitions each, player_id key)", "", RGBColor(0xE8, 0xE0, 0xF7), text_color=PURPLE, font_size_title=10)

add_text_box(slide, 0.5, 6.4, 12, 0.6,
    "Data flow: Player input (WebSocket) -> Game Server (physics + feature extraction) -> Kafka (4 topics) -> ML Service (inference) -> Kafka (alerts) -> Grafana (dashboards)",
    font_size=12, color=MUTED)


# ===== SLIDE 4: DATA PIPELINE & FEATURE ENGINEERING =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Data Pipeline & Feature Engineering", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "From raw input to computed features via Kafka streaming", font_size=16, color=MUTED)

add_text_box(slide, 0.8, 1.6, 3.8, 0.4, "Telemetry Generation (60Hz)", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 0.8, 2.0, 3.8, 2.2, [
    "Player sends raw input (keys, mouse) via WebSocket",
    "Server computes position, velocity, aim, hit detection",
    "Spatial context: nearest enemy distance, angle, LOS",
    "18 raw fields per player per tick to Kafka",
], font_size=13)

add_text_box(slide, 5, 1.6, 3.8, 0.4, "Feature Engineering (Go)", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 5, 2.0, 3.8, 2.2, [
    "3 sliding windows: 1s (60 ticks), 5s (300), 30s (1800)",
    "18 computed features: aim dynamics, combat stats, movement, spatial correlation",
    "Ring buffer per player, 60s cleanup for inactive",
    "Feature vectors sent to ML on every tick",
], font_size=13)

add_text_box(slide, 9.2, 1.6, 3.8, 0.4, "Kafka Topics (KRaft)", font_size=18, bold=True, color=TEXT)
add_table(slide, 9.2, 2.0, 3.5, 1.6,
    ["Topic", "Purpose"],
    [
        ["telemetry.raw", "Raw state/tick"],
        ["events.kills", "Kill events"],
        ["features.computed", "18-feature vectors"],
        ["alerts.detections", "Rule + ML alerts"],
    ], font_size=10)

add_text_box(slide, 0.8, 4.5, 12, 0.4, "Cheat Simulation (Server-Side)", font_size=18, bold=True, color=TEXT)
add_table(slide, 0.8, 5.0, 11.5, 2,
    ["Cheat", "Toggle", "Server-Side Behavior", "Detectable Signature", "Rule Threshold"],
    [
        ["Aimbot", "F5", "Snaps aim to nearest enemy via atan2()", "aim_delta > 2.0 rad/tick", "Normal: 0-0.5 rad"],
        ["Speedhack", "F6", "2.5x velocity multiplier", "speed > 7.0 (cap is 5.0)", "Normal: 0-5.0"],
        ["Wallhack", "F7", "Aim tracks enemy through walls", "prefire_ratio > 60%", "Normal: 0-15%"],
        ["Triggerbot", "F8", "Auto-fires within 0.1 rad of enemy", "reaction_time < 3 ticks", "Normal: 9-15 ticks"],
    ], font_size=11)


# ===== SLIDE 5: RULE ENGINE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Detection Layer 1: Rule Engine", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "Physics-derived thresholds with zero false positive rate", font_size=16, color=MUTED)

add_table(slide, 0.5, 1.7, 12.3, 3,
    ["Rule", "Condition", "Detects", "Normal Range", "Cheat Range", "Rationale"],
    [
        ["speed_cap", "speed_max_1s > 7.0", "Speedhack", "0 - 5.0", "~12.5", "Physics cap 5.0; 40% buffer"],
        ["aim_snap", "aim_delta_max_1s > 2.0", "Aimbot", "0 - 0.5 rad", "~3.14 rad", "Mouse can't exceed 0.5 rad/tick"],
        ["inhuman_accuracy", "hit_rate_5s > 85%, shots > 30", "Aimbot", "20 - 50%", "~100%", "85% over 30+ shots unsustainable"],
        ["aim_lock", "aim_lock_ratio > 90%", "Aimbot", "5 - 30%", "~100%", "Aim within 0.1 rad 90%+ of time"],
        ["prefire", "prefire_ratio > 60%, shots > 20", "Wallhack", "0 - 15%", "> 60%", "Consistently shooting through walls"],
        ["triggerbot_reaction", "reaction_time < 3 ticks", "Triggerbot", "9 - 15 ticks", "1 - 2 ticks", "3 ticks = 50ms; human min ~150ms"],
    ], font_size=11)

add_text_box(slide, 0.8, 5.3, 11.5, 0.8,
    "Key insight: Every threshold sits in a dead zone between what's physically possible for a human and what the cheat "
    "produces. No statistical tuning needed. Derived directly from game physics constants.",
    font_size=15, bold=True, color=TEXT)

add_bullet_list(slide, 0.8, 6.2, 11.5, 1, [
    "Rule alerts include specific cheat type classification (aimbot, speedhack, wallhack, triggerbot)",
    "5-second per-player cooldown prevents alert flooding while detection continues",
], font_size=13)


# ===== SLIDE 6: ML MODELS (XGBoost + LSTM) =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Detection Layers 2 & 3: ML Models", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "XGBoost (supervised, known cheats) + LSTM Autoencoder (unsupervised, novel cheats)", font_size=16, color=MUTED)

add_text_box(slide, 0.8, 1.6, 5.8, 0.4, "XGBoost Binary Classifier", font_size=20, bold=True, color=ACCENT)
add_bullet_list(slide, 0.8, 2.0, 5.8, 1.8, [
    "Input: 18-feature normalized vector (per tick)",
    "Output: Binary probability (none vs cheater), threshold > 80%",
    "Explainability: Top-3 feature importances per alert",
    "Split by player ID (80/20) to prevent data leakage",
], font_size=13)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(5.8), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "max_depth: 6 | learning_rate: 0.1 | n_estimators: 200 | subsample: 0.8 | objective: binary:logistic"
run.font.size = Pt(11)
run.font.name = "Consolas"
run.font.color.rgb = MUTED

add_text_box(slide, 0.8, 4.7, 5.8, 0.4, "Why XGBoost?", font_size=16, bold=True, color=TEXT)
add_bullet_list(slide, 0.8, 5.1, 5.8, 1.8, [
    "Outperforms neural nets on tabular data",
    "Microsecond inference (real-time capable)",
    "Built-in feature importance for explainability",
], font_size=13)

add_text_box(slide, 7, 1.6, 5.8, 0.4, "LSTM Autoencoder (Anomaly Detection)", font_size=20, bold=True, color=ACCENT)
add_bullet_list(slide, 7, 2.0, 5.8, 1.8, [
    "Encoder: LSTM(18->64) then LSTM(64->32)",
    "Decoder: RepeatVector(60) -> LSTM(32->32) -> LSTM(32->64) -> Dense(18)",
    "Input: 60 consecutive feature vectors (1 second at 60Hz)",
    "Loss: MSE reconstruction error; threshold: mean + 3 std",
], font_size=13)

add_text_box(slide, 7, 3.9, 5.8, 0.4, "How It Works", font_size=16, bold=True, color=TEXT)
add_bullet_list(slide, 7, 4.3, 5.8, 1, [
    "Trained only on clean (non-cheating) player data",
    "Learns to reconstruct \"normal\" behavior sequences",
    "Cheating = high reconstruction error (anomaly)",
], font_size=13)

add_text_box(slide, 7, 5.5, 5.8, 0.4, "Why Unsupervised?", font_size=16, bold=True, color=TEXT)
add_text_box(slide, 7, 5.9, 5.8, 1.2,
    "XGBoost only detects cheats it was trained on. A novel cheat type "
    "not in training data would be classified as \"none\". The autoencoder catches it "
    "because the behavior deviates from normal, regardless of cheat type.",
    font_size=13, color=MUTED)


# ===== SLIDE 7: TRAINING DATASET =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Training: CS2CD Dataset + Synthetic Data", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "Public Counter-Strike 2 Cheat Detection dataset with VAC-confirmed labels, augmented with game data", font_size=16, color=MUTED)

add_text_box(slide, 0.8, 1.6, 3.8, 0.4, "CS2CD Public Dataset", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 0.8, 2.0, 3.8, 2.2, [
    "795 real CS2 competitive matches",
    "451 with VAC-banned cheaters, 344 clean",
    "Ground-truth labels from Valve's VAC system",
    "Parquet format, ~247 columns per tick",
], font_size=13)

add_text_box(slide, 5, 1.6, 3.8, 0.4, "Our Feature Mapping", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 5, 2.0, 3.8, 2.2, [
    "Map CS2 3D fields to our 2D schema (drop Z axis)",
    "Recompute spatial context from positions",
    "Extract same 18 sliding-window features",
    "Tick downsampling (30K max) for memory efficiency",
], font_size=13)

add_text_box(slide, 9.2, 1.6, 3.8, 0.4, "Synthetic Augmentation", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 9.2, 2.0, 3.8, 2.2, [
    "Bot-generated labeled telemetry from our game",
    "Bridges 3D/2D feature distribution gap",
    "Bots cycle all 4 cheat types randomly",
    "Exact feature distributions our game produces",
], font_size=13)

add_text_box(slide, 0.8, 4.5, 12, 0.4, "Alert Combiner (Ensemble Logic)", font_size=18, bold=True, color=TEXT)
add_table(slide, 0.8, 5.0, 5.5, 1.8,
    ["Scenario", "Source", "Confidence"],
    [
        ["XGBoost only", "xgboost", "Model probability"],
        ["Autoencoder only", "autoencoder", "Anomaly score (0-1)"],
        ["Both trigger", "ensemble", "max(both)"],
        ["Rule engine", "rule-engine", "1.0 (deterministic)"],
    ], font_size=11)

add_text_box(slide, 7, 4.5, 5.5, 0.4, "Why Three Detection Layers?", font_size=18, bold=True, color=TEXT)
add_bullet_list(slide, 7, 5.0, 5.5, 2, [
    "Rules: Instant, deterministic, zero false positives for blatant cheats",
    "XGBoost: Catches subtle multi-feature combinations individually within normal ranges",
    "LSTM: Catches novel cheats via behavioral deviation from normal patterns",
    "Graceful degradation: rules still work if ML service is down",
], font_size=13)


# ===== SLIDE 8: OBSERVABILITY =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Observability & Dashboards", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "Prometheus (17 metrics) + Grafana (2 dashboards, auto-provisioned via Docker)", font_size=16, color=MUTED)

add_text_box(slide, 0.8, 1.6, 5.5, 0.4, "Pipeline Health Dashboard", font_size=20, bold=True, color=TEXT)
add_bullet_list(slide, 0.8, 2.0, 5.5, 3, [
    "Active Players (gauge)",
    "Tick Duration histogram (16.6ms budget at 60Hz)",
    "WebSocket Connections",
    "Telemetry / Feature publish rates",
    "Kafka error counter",
    "Feature processing latency (p99)",
    "Model loaded status (XGBoost + Autoencoder)",
], font_size=14)

add_text_box(slide, 7, 1.6, 5.5, 0.4, "Detection Analytics Dashboard", font_size=20, bold=True, color=TEXT)
add_bullet_list(slide, 7, 2.0, 5.5, 3, [
    "Rule Engine Alerts (stacked by rule name)",
    "ML Alerts by Model (xgboost / autoencoder / ensemble)",
    "Total Alert Rate (color-coded green/yellow/red)",
    "Predictions by Type (pie chart: none vs cheater)",
    "Anomaly Score Distribution (heatmap)",
    "XGBoost + Autoencoder inference latency",
    "Training samples collected, service mode",
], font_size=14)

add_text_box(slide, 0.8, 5.5, 11.5, 1.2,
    "Both dashboards are auto-provisioned via Docker volume mounts. No manual setup required. "
    "Prometheus scrapes game server and ML service every 5 seconds. Grafana provides real-time "
    "visualization of pipeline health and detection analytics.",
    font_size=14, color=MUTED)


# ===== SLIDE 9: LIVE DEMO =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN)
add_text_box(slide, 1.5, 1.0, 10, 1, "Live Demo", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2, 2.2, 9, 0.6, "Demonstrating real-time cheat detection end-to-end", font_size=20, color=RGBColor(0xC8, 0xE6, 0xC9), alignment=PP_ALIGN.CENTER)

add_bullet_list(slide, 2, 3.2, 9, 3.5, [
    "Step 1: Open game client at localhost:80, show clean gameplay with bots",
    "Step 2: Open Grafana dashboards (localhost:3000), show zero alerts baseline",
    "Step 3: Toggle Aimbot (F5), watch aim_snap + aim_lock rules fire instantly on dashboard",
    "Step 4: Toggle Speedhack (F6), speed_cap rule fires, ML prediction changes to \"cheater\"",
    "Step 5: Toggle Wallhack (F7), prefire rule fires after shot accumulation",
    "Step 6: Toggle Triggerbot (F8), reaction_time rule fires within 1 second",
    "Step 7: Disable all cheats (F5-F8 again), alerts stop, pie chart returns to \"none\"",
    "Step 8: Show ML service health endpoint confirming inference mode with trained models",
], font_size=15, color=WHITE, bold_prefix=True)


# ===== SLIDE 10: EVALUATION METRICS =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Evaluation & Results", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "Measurable performance across latency, throughput, and accuracy", font_size=16, color=MUTED)

add_metric_box(slide, 0.8, 1.7, "< 17ms", "Rule Engine Latency\n(same-tick detection)")
add_metric_box(slide, 3.8, 1.7, "< 5s", "End-to-End Alert\n(cheat to dashboard)")
add_metric_box(slide, 6.8, 1.7, "540/s", "Telemetry Throughput\n(events per second)")
add_metric_box(slide, 9.8, 1.7, "0%", "Rule Engine FP Rate\n(physics thresholds)")

add_metric_box(slide, 0.8, 3.6, "18", "Computed Features\n(sliding window)")
add_metric_box(slide, 3.8, 3.6, "41K", "Training Samples\n(CS2CD + synthetic)")
add_metric_box(slide, 6.8, 3.6, "7", "Docker Containers\n(fully orchestrated)")
add_metric_box(slide, 9.8, 3.6, "17", "Prometheus Metrics\n(2 dashboards)")

add_text_box(slide, 0.8, 5.5, 5.5, 0.4, "XGBoost Results", font_size=16, bold=True, color=TEXT)
add_text_box(slide, 0.8, 5.9, 5.5, 1,
    "Trained on 31K CS2CD + 10K synthetic samples. Binary classification (none vs cheater). "
    "Cross-domain accuracy bounded by 3D/2D feature distribution gap, but correctly classifies cheats in live gameplay.",
    font_size=13, color=MUTED)

add_text_box(slide, 7, 5.5, 5.5, 0.4, "LSTM Autoencoder Results", font_size=16, bold=True, color=TEXT)
add_text_box(slide, 7, 5.9, 5.5, 1,
    "Trained on clean player sequences only. Anomaly threshold calibrated at mean + 3 std of reconstruction errors. "
    "Fires on behavioral patterns unseen during training, complementing supervised detection.",
    font_size=13, color=MUTED)


# ===== SLIDE 11: LIMITATIONS & FUTURE WORK =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.6, "Limitations & Future Work", font_size=32, bold=True, color=TEXT)
add_text_box(slide, 0.8, 1.0, 10, 0.4, "Honest assessment of constraints and potential improvements", font_size=16, color=MUTED)

add_text_box(slide, 0.8, 1.6, 5.5, 0.4, "Current Limitations", font_size=20, bold=True, color=TEXT)
add_bullet_list(slide, 0.8, 2.0, 5.5, 3, [
    "Feature distribution gap: CS2 (3D) vs our game (2D) bounds cross-domain accuracy",
    "Scale: Tested with 9 players; horizontal scaling (consumer groups) not implemented",
    "Information-only cheats: Wallhack without aim tracking produces subtle signals",
    "No online learning: Models require manual retrain when player behavior drifts",
    "CPU-only inference: LSTM autoencoder limited throughput without GPU",
], font_size=14)

add_text_box(slide, 7, 1.6, 5.5, 0.4, "Future Improvements", font_size=20, bold=True, color=TEXT)
add_bullet_list(slide, 7, 2.0, 5.5, 3, [
    "Temporal transformer: Better long-range pattern detection than LSTM",
    "Graph neural network: Multi-player collusion detection via kill-graph",
    "Adversarial training: Generator creates stealthy cheats, detector learns",
    "Federated learning: Train across multiple game titles without sharing data",
    "Online threshold calibration: Autoencoder threshold adapts to population shifts",
    "Tiered inference: Run expensive LSTM only for XGBoost/rule-flagged players",
], font_size=14)

add_text_box(slide, 0.8, 5.5, 11.5, 0.4, "Key Contributions", font_size=20, bold=True, color=TEXT)
add_bullet_list(slide, 0.8, 5.9, 11.5, 1.2, [
    "End-to-end real-time pipeline: from 60Hz telemetry to Grafana alert in under 5 seconds",
    "Dual-paradigm ML: supervised XGBoost for known cheats + unsupervised LSTM for novel cheats",
    "Physics-derived rule engine: zero false positives, no ML dependency, deterministic detection",
    "Trained on public CS2CD dataset (795 real matches) for reproducibility and academic rigor",
], font_size=13)


# ===== SLIDE 12: THANK YOU =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, ACCENT)
add_text_box(slide, 1.5, 2.0, 10, 1.5, "Thank You", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2, 4.0, 9, 1, "Questions?", font_size=28, color=RGBColor(0xC5, 0xD2, 0xF0), alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2, 5.5, 9, 1,
    "Source Code: Docker Compose one-command deployment\n"
    "Tech Stack: Go 1.22 | Python 3.11 | Apache Kafka | XGBoost | LSTM | Prometheus | Grafana",
    font_size=14, color=RGBColor(0xA5, 0xB2, 0xE0), alignment=PP_ALIGN.CENTER)


# SAVE
output_path = "/mnt/c/Users/reaul/Downloads/UOIT/Spring 2026/Real-Time Data Analytics for IoT/Project/cheat-detection/docs/final-presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
